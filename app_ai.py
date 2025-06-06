# Paydo AI PPT 생성기 with KoSimCSE + KSS 의미 단위 분할 적용

import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io
import re
import textwrap
import docx
from io import BytesIO
from sentence_transformers import SentenceTransformer, util
import kss

# Streamlit 세팅
st.set_page_config(page_title="Paydo AI PPT", layout="centered")
st.title("🎬 AI PPT 생성기 (KoSimCSE + 의미 단위 분할)")

@st.cache_resource
def load_model():
    return SentenceTransformer("jhgan/ko-sbert-nli")

model = load_model()

# Word 파일 텍스트 추출
def extract_text_from_word(uploaded_file):
    try:
        file_bytes = BytesIO(uploaded_file.read())
        doc = docx.Document(file_bytes)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        if not paragraphs:
            raise ValueError("Word 파일이 비어 있습니다.")  # 파일이 비어있는 경우 예외 발생
        return paragraphs
    except Exception as e:
        st.error(f"Word 파일 처리 오류: {e}")
        return None

# 텍스트 줄 수 계산
def calculate_text_lines(text, max_chars_per_line):
    lines = 0
    paragraphs = text.split('\n')
    for paragraph in paragraphs:
        if not paragraph:
            lines += 1
        else:
            lines += len(textwrap.wrap(paragraph, width=max_chars_per_line, break_long_words=True))
    return lines

# 의미 단위 기준 문장 분할 (KSS 사용)
def smart_sentence_split(text):
    return kss.split_sentences(text)

# 불완전 문장 판단 및 병합
connective_pattern = re.compile(r'^(그리고|하지만|그러나|또한|그래서|즉|또|그러면|그런데)$')

def is_incomplete(sentence):
    return (
        sentence.endswith(('은', '는', '이', '가', '을', '를', '에', '으로', '고', '와', '과'))
        or len(sentence.strip()) < 8
        or bool(connective_pattern.match(sentence.strip()))
    )

# 슬라이드 분할 with 의미 단위 + 문맥 유사도
def split_text_into_slides_with_similarity(text_paragraphs, max_lines_per_slide, max_chars_per_line_ppt, similarity_threshold=0.85):
    slides, split_flags = [], []
    slide_number = 1
    current_text = ""
    current_lines = 0
    needs_check = False

    for paragraph in text_paragraphs:
        sentences = smart_sentence_split(paragraph)
        if not sentences:
            continue

        merged_sentences = []
        buffer = ""
        for sentence in sentences:
            if buffer:
                buffer += " " + sentence
                if not is_incomplete(sentence):
                    merged_sentences.append(buffer.strip())
                    buffer = ""
            else:
                if is_incomplete(sentence):
                    buffer = sentence
                else:
                    merged_sentences.append(sentence)
        if buffer:
            merged_sentences.append(buffer.strip())

        embeddings = model.encode(merged_sentences)

        i = 0
        while i < len(merged_sentences):
            sentence = merged_sentences[i]
            sentence_lines = calculate_text_lines(sentence, max_chars_per_line_ppt)

            if sentence_lines <= 2 and i + 1 < len(merged_sentences):
                next_sentence = merged_sentences[i + 1]
                merged = sentence + " " + next_sentence
                merged_lines = calculate_text_lines(merged, max_chars_per_line_ppt)
                if merged_lines <= max_lines_per_slide:
                    sentence = merged
                    sentence_lines = merged_lines
                    i += 1

            if sentence_lines > max_lines_per_slide:
                wrapped_lines = textwrap.wrap(sentence, width=max_chars_per_line_ppt, break_long_words=True)
                temp_text = ""
                temp_lines = 0
                for line in wrapped_lines:
                    line_lines = calculate_text_lines(line, max_chars_per_line_ppt)
                    if temp_lines + line_lines <= max_lines_per_slide:
                        temp_text += line + "\n"
                        temp_lines += line_lines
                    else:
                        slides.append(temp_text.strip())
                        split_flags.append(True)
                        slide_number += 1
                        temp_text = line + "\n"
                        temp_lines = line_lines
                if temp_text:
                    slides.append(temp_text.strip())
                    split_flags.append(True)
                    slide_number += 1
                current_text = ""
                current_lines = 0
                i += 1
                continue

            if current_lines + sentence_lines <= max_lines_per_slide:
                current_text += sentence + "\n"
                current_lines += sentence_lines
            else:
                slides.append(current_text.strip())
                split_flags.append(needs_check)
                slide_number += 1
                current_text = sentence + "\n"
                current_lines = sentence_lines
                needs_check = False
            i += 1

    if current_text:
        slides.append(current_text.strip())
        split_flags.append(needs_check)

    return slides, split_flags

# [create_ppt, add_text_to_slide, add_check_needed_shape, add_end_mark 함수는 이전 정의대로 유지]
def create_ppt(slide_texts, split_flags, max_chars_per_line_in_ppt=18, font_size=54):
    """슬라이드 텍스트를 기반으로 PPT를 생성하고, '확인 필요!' 표시 등을 추가합니다."""

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    total_slides = len(slide_texts)

    for i, text in enumerate(slide_texts):
        try:
            logging.debug(f"슬라이드 {i+1}에 텍스트 추가")
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_text_to_slide(slide, text, font_size, PP_ALIGN.CENTER, max_chars_per_line_in_ppt)
            if split_flags[i]:
                add_check_needed_shape(slide)  # 슬라이드 번호 인자 제거
            if i == total_slides - 1:
                add_end_mark(slide)
        except Exception as e:
            st.error(f"오류: 슬라이드 생성 실패 (슬라이드 {i+1}): {e}")
            return None

    return prs

def add_text_to_slide(slide, text, font_size, alignment, max_chars_per_line):
    """슬라이드에 텍스트를 추가하고, 폰트, 크기, 정렬 등을 설정합니다."""

    try:
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(6.2))
        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        text_frame.word_wrap = True

        wrapped_lines = textwrap.wrap(text, width=max_chars_per_line, break_long_words=True)
        for line in wrapped_lines:
            p = text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(font_size)
            p.font.name = 'Noto Color Emoji'
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = alignment
            p.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

        text_frame.auto_size = None
        logging.debug(f"텍스트 추가됨")
    except Exception as e:
        st.error(f"오류: 슬라이드에 텍스트 추가 중 오류 발생: {e}")
        raise

def add_end_mark(slide):
    """마지막 슬라이드에 '끝' 표시를 추가합니다."""

    end_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(10),
        Inches(6),
        Inches(2),
        Inches(1)
    )
    end_shape.fill.solid()
    end_shape.fill.fore_color.rgb = RGBColor(255, 0, 0)
    end_shape.line.color.rgb = RGBColor(0, 0, 0)

    end_text_frame = end_shape.text_frame
    end_text_frame.clear()
    p = end_text_frame.paragraphs[0]
    p.text = "끝"
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(255, 255, 255)
    end_text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p.alignment = PP_ALIGN.CENTER

def add_check_needed_shape(slide):
    """확인 필요한 슬라이드에 '확인 필요!' 상자를 추가합니다."""

    check_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(0.3),
        Inches(2.5),
        Inches(0.5)
    )
    check_shape.fill.solid()
    check_shape.fill.fore_color.rgb = RGBColor(255, 255, 0)
    check_shape.line.color.rgb = RGBColor(0, 0, 0)

    check_text_frame = check_shape.text_frame
    check_text_frame.clear()
    p = check_text_frame.paragraphs[0]
    p.text = "확인 필요!"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    text_frame = check_shape.text_frame
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p.alignment = PP_ALIGN.CENTER

# UI 입력
uploaded_file = st.file_uploader("📄 Word 파일 업로드", type=["docx"])
text_input = st.text_area("또는 텍스트 직접 입력:", height=300)

max_lines = st.slider("슬라이드당 최대 줄 수", 1, 10, 4)
max_chars = st.slider("한 줄당 최대 글자 수", 10, 100, 18)
font_size = st.slider("폰트 크기", 10, 60, 54)
sim_threshold = st.slider("문맥 유사도 기준", 0.0, 1.0, 0.85, step=0.05)

if st.button("🚀 PPT 생성"):
    paragraphs = []
    if uploaded_file:
        paragraphs = extract_text_from_word(uploaded_file)
    elif text_input.strip():
        paragraphs = [p.strip() for p in text_input.split("\n\n") if p.strip()]
    else:
        st.warning("Word 파일을 업로드하거나 텍스트를 입력하세요.")
        st.stop()

    if not paragraphs:
        st.error("유효한 텍스트가 없습니다.")
        st.stop()

    with st.spinner("PPT 생성 중..."):
        slides, flags = split_text_into_slides_with_similarity(
            paragraphs, max_lines, max_chars, sim_threshold
        )
        ppt = create_ppt(slides, flags, max_chars, font_size)

        if ppt:
            ppt_io = io.BytesIO()
            ppt.save(ppt_io)
            ppt_io.seek(0)
            st.download_button("📥 PPT 다운로드", ppt_io, "paydo_script_ai.pptx",
                               mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
            st.success(f"총 {len(slides)}개의 슬라이드가 생성되었습니다.")
            if any(flags):
                flagged = [i+1 for i, f in enumerate(flags) if f]
                st.warning(f"⚠️ 확인이 필요한 슬라이드: {flagged}")