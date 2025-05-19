import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io
import re
import textwrap
import docx
from datetime import datetime

# Word 파일에서 텍스트 추출하는 함수 (기존 코드와 동일)
def extract_text_from_word(file_like_object):
    """업로드된 파일 객체에서 텍스트를 추출합니다."""
    doc = docx.Document(file_like_object)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

# 문장이 차지할 줄 수 계산 (기존 코드와 동일)
def calculate_text_lines(text, max_chars_per_line):
    lines = 0
    paragraphs = text.split('\n')
    for paragraph in paragraphs:
        if not paragraph:
            lines += 1
        else:
            lines += len(textwrap.wrap(paragraph, width=max_chars_per_line, break_long_words=True))
    return lines

# 텍스트를 슬라이드로 분할 및 그룹화 (기존 코드와 동일)
def split_and_group_text(text, max_lines_per_slide, max_chars_per_line_ppt):
    slides = []
    split_flags = []
    lines = text.strip().split('\n')

    for line in lines:
        line = line.strip()
        line_count = calculate_text_lines(line, max_chars_per_line_ppt)

        if not slides:
            slides.append(line)
            split_flags.append(False)
        elif calculate_text_lines(slides[-1] + "\n" + line, max_chars_per_line_ppt) <= max_lines_per_slide:
            slides[-1] += "\n" + line
            split_flags[-1] = False
        else:
            slides.append(line)
            split_flags.append(False)

    final_slides = []
    final_split_flags = []
    max_chars_per_segment = 60

    for i, slide_text in enumerate(slides):
        if calculate_text_lines(slide_text, max_chars_per_line_ppt) > max_lines_per_slide:
            original_sentence = slide_text.replace('\n', ' ')
            sub_sentences = re.split(r'(?<=[.?!;])\s+', original_sentence.strip())
            temp_slide_text = ""
            temp_slide_lines = 0
            is_forced_split = False
            for sub_sentence in sub_sentences:
                sub_sentence = sub_sentence.strip()
                sub_sentence_lines = calculate_text_lines(sub_sentence, max_chars_per_line_ppt)
                if temp_slide_lines + sub_sentence_lines <= max_lines_per_slide:
                    if temp_slide_text:
                        temp_slide_text += " "
                    temp_slide_text += sub_sentence
                    temp_slide_lines += sub_sentence_lines
                else:
                    final_slides.append(temp_slide_text)
                    final_split_flags.append(is_forced_split)
                    temp_slide_text = sub_sentence
                    temp_slide_lines = sub_sentence_lines
                    is_forced_split = False

            if temp_slide_text:
                if calculate_text_lines(temp_slide_text, max_chars_per_line_ppt) > max_lines_per_slide:
                    words = temp_slide_text.split()
                    segment = ""
                    for word in words:
                        if len(segment.replace(" ", "")) + len(word) + (1 if segment else 0) <= max_chars_per_segment:
                            if segment:
                                segment += " "
                            segment += word
                        else:
                            final_slides.append(segment)
                            final_split_flags.append(True)
                            segment = word
                            is_forced_split = True
                    if segment:
                        final_slides.append(segment)
                        final_split_flags.append(True)
                else:
                    final_slides.append(temp_slide_text)
                    final_split_flags.append(False)
        else:
            final_slides.append(slide_text)
            final_split_flags.append(False)

    final_slides = [slide for slide in final_slides if slide.strip()]
    final_split_flags = final_split_flags[:len(final_slides)]

    return final_slides, final_split_flags

# PPT 생성 함수 (기존 코드와 동일)
def create_ppt(slide_texts, split_flags, max_chars_per_line_in_ppt=18, font_size=54):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    total_slides = len(slide_texts)

    for i, text in enumerate(slide_texts):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_text_to_slide(slide, text, font_size, PP_ALIGN.CENTER)
        add_slide_number(slide, i + 1, total_slides)
        if split_flags[i]:
            add_check_needed_shape(slide)
        if i == total_slides - 1:
            add_end_mark(slide)

    return prs

# 슬라이드에 텍스트 추가 (기존 코드와 동일)
def add_text_to_slide(slide, text, font_size, alignment):
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(6.2))
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP  # 상단 정렬 명시적으로 설정
    text_frame.word_wrap = True

    wrapped_lines = textwrap.wrap(text, width=18, break_long_words=True)  # 긴 단어 분리 활성화
    text_frame.clear()
    for line in wrapped_lines:
        p = text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = 'Noto Color Emoji'
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = alignment
        p.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    # 텍스트 박스의 자동 맞춤 기능 제거 (상단 정렬에 영향 줄 수 있음)
    text_frame.auto_size = None
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

# 슬라이드 번호 추가 (기존 코드와 동일)
def add_slide_number(slide, current, total):
    footer_box = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4))
    footer_text_frame = footer_box.text_frame
    footer_text_frame.clear()
    p = footer_text_frame.paragraphs[0]
    p.text = f"{current} / {total}"
    p.font.size = Pt(18)
    p.font.name = '맑은 고딕'
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.RIGHT

# '끝' 모양 추가 (기존 코드와 동일)
def add_end_mark(slide):
    end_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(10),
        Inches(6),
        Inches(2),
        Inches(1)
    )
    end_shape.fill.solid()
    end_shape.fill.fore_color.rgb = RGBColor(255, 0, 0)
    end_shape.line.color.rgb = RGBColor(0, 0, 0)

    end_text_frame = end_shape.text_frame
    end_text_frame.clear()
    p = end_text_frame.paragraphs[0]
    p.text = "끝"
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(255, 255, 255)
    end_text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p.alignment = PP_ALIGN.CENTER

# '확인 필요!' 모양 추가 (기존 코드와 동일)
def add_check_needed_shape(slide):
    check_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(0.3),
        Inches(2),
        Inches(0.5)
    )
    check_shape.fill.solid()
    check_shape.fill.fore_color.rgb = RGBColor(255, 255, 0)
    check_shape.line.color.rgb = RGBColor(0, 0, 0)

    check_text_frame = check_shape.text_frame
    check_text_frame.clear()
    p = check_text_frame.paragraphs[0]
    p.text = "확인 필요!"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    text_frame = check_shape.text_frame
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p.alignment = PP_ALIGN.CENTER

# Streamlit UI
st.set_page_config(page_title="Paydo", layout="centered")
st.title("🎬 Paydo 촬영 대본 PPT 자동 생성기")

# 사이드바 설정 (기존 코드와 동일)
with st.sidebar:
    st.header("⚙️ PPT 설정")
    max_lines_per_slide_input = st.slider(
        "📄 슬라이드당 최대 줄 수:", min_value=1, max_value=10, value=5, key="max_lines_slider"
    )
    st.caption("한 슬라이드에 들어갈 최대 줄 수를 설정합니다.")
    max_chars_per_line_ppt_input = st.slider(
        "📏 한 줄당 최대 글자 수 (PPT 표시):", min_value=3, max_value=30, value=18, key="max_chars_slider_ppt"
    )
    st.caption("PPT에 표시될 텍스트의 한 줄당 최대 글자 수를 설정합니다.")
    font_size_input = st.slider(
        "🅰️ 폰트 크기:", min_value=10, max_value=60, value=54, key="font_size_slider"
    )
    st.caption("PPT 텍스트의 폰트 크기를 설정합니다.")

# 메인 화면 디자인 개선 (기존 코드와 동일)
with st.container():
    st.markdown("### 📝 촬영 대본 입력")
    st.markdown(
        """
    Word 파일(.docx)을 업로드하거나, 텍스트를 직접 입력하세요.
    """
    )
    with st.form(key="input_form"):  # 입력 영역을 form으로 묶음
        col1, col2 = st.columns(2)  # 2개의 컬럼으로 나눔
        with col1:
            uploaded_file = st.file_uploader(
                "Word 파일 업로드", type=["docx"], help="docx 형식의 파일만 지원됩니다."
            )
        with col2:
            text_input = st.text_area(
                "텍스트 직접 입력",
                height=200,
                placeholder="여기에 텍스트를 입력하세요...",
                help="텍스트를 직접 입력할 수 있습니다.",
            )
        submit_button = st.form_submit_button("🚀 PPT 만들기")  # key 인자 제거

if submit_button:  # 버튼이 눌렸을 때만 처리
    text = ""
    from io import BytesIO  # 파일 상단에 이미 import 되어 있다면 생략

if uploaded_file is not None:
    try:
        file_bytes = BytesIO(uploaded_file.read())  # 핵심: 강제 래핑
        text = extract_text_from_word(file_bytes)
    except Exception as e:
        st.error(f"📄 파일을 읽는 중 오류가 발생했습니다: {e}")
        st.stop()

elif text_input.strip():
        text = text_input
else:
        st.error("Word 파일을 업로드하거나 텍스트를 입력하세요.")
        st.stop()

    # 파일 제목 설정 (수정됨)
now = datetime.now()
date_string = now.strftime("%y%m%d")  # YYMMDD 형식
ppt_filename = f"[촬영 대본] paydo_script_{date_string}.pptx"  # 파일 이름 통일

    # PPT 생성 진행 표시 (기존 코드와 동일)
with st.spinner("PPT 생성 중..."):
        slide_texts, split_flags = split_and_group_text(
            text,
            max_lines_per_slide=max_lines_per_slide_input,
            max_chars_per_line_ppt=max_chars_per_line_ppt_input,
        )
        ppt = create_ppt(
            slide_texts,
            split_flags,
            max_chars_per_line_in_ppt=max_chars_per_line_ppt_input,
            font_size=font_size_input,
        )

if ppt:
        ppt_io = io.BytesIO()
        ppt.save(ppt_io)
        ppt_io.seek(0)

        st.success("PPT 생성 완료! 아래 버튼을 눌러 다운로드하세요.")
        st.download_button(
            label="📥 PPT 다운로드",
            data=ppt_io,
            file_name=ppt_filename,  # 동적으로 생성된 파일 이름 사용
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key="download_button"
        )
        if any(split_flags):
            split_slide_numbers = [i + 1 for i, flag in enumerate(split_flags) if flag]
            st.warning(
                f"❗️ 일부 슬라이드({split_slide_numbers})는 한 문장이 너무 길어 분할되었습니다. PPT를 확인하여 가독성을 검토해주세요."
            )
else:
        st.error("❌ PPT 생성에 실패했습니다.")